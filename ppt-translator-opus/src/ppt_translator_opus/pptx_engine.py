"""PPTX text extraction and translation apply-back engine."""

import json
import copy
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.text import MSO_AUTO_SIZE
from lxml import etree


# ── Namespaces for SmartArt / raw XML ──────────────────────────────

NSMAP = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "dgm": "http://schemas.openxmlformats.org/drawingml/2006/diagram",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}


# ── Extract helpers ────────────────────────────────────────────────

def _extract_runs(paragraph) -> list[dict]:
    """Extract run-level data from a paragraph."""
    runs = []
    for i, run in enumerate(paragraph.runs):
        text = run.text
        if not text:
            continue
        run_data = {
            "index": i,
            "text": text,
        }
        # Record format info for reference (not used by LLM, but useful for debugging)
        if run.font.bold is not None:
            run_data["bold"] = run.font.bold
        if run.font.italic is not None:
            run_data["italic"] = run.font.italic
        if run.font.size is not None:
            run_data["font_size"] = run.font.size.pt
        runs.append(run_data)
    return runs


def _extract_paragraphs(text_frame) -> list[dict]:
    """Extract all runs from all paragraphs in a text frame, flattened with global index."""
    all_runs = []
    global_idx = 0
    for para in text_frame.paragraphs:
        for run in para.runs:
            text = run.text
            if not text:
                continue
            run_data = {
                "index": global_idx,
                "text": text,
            }
            if run.font.bold is not None:
                run_data["bold"] = run.font.bold
            if run.font.italic is not None:
                run_data["italic"] = run.font.italic
            if run.font.size is not None:
                run_data["font_size"] = run.font.size.pt
            all_runs.append(run_data)
            global_idx += 1
    return all_runs


def _extract_shape(shape, slide_num: int, counter: dict) -> list[dict]:
    """Extract text elements from a single shape, recursively handling groups."""
    elements = []

    # Group shapes: recurse into children
    if shape.shape_type is not None and shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
        try:
            for child_shape in shape.shapes:
                elements.extend(_extract_shape(child_shape, slide_num, counter))
        except Exception:
            pass
        return elements

    # Table
    if shape.has_table:
        table = shape.table
        for row_idx, row in enumerate(table.rows):
            for col_idx, cell in enumerate(row.cells):
                if cell.text_frame:
                    runs = _extract_paragraphs(cell.text_frame)
                    if runs:
                        counter["idx"] += 1
                        elements.append({
                            "id": f"s{slide_num}_table_{counter['idx']}_r{row_idx}c{col_idx}",
                            "type": "table_cell",
                            "runs": runs,
                        })
        return elements

    # Text frame (text boxes, titles, subtitles, etc.)
    if shape.has_text_frame:
        runs = _extract_paragraphs(shape.text_frame)
        if runs:
            counter["idx"] += 1
            elements.append({
                "id": f"s{slide_num}_shape_{counter['idx']}",
                "type": "textbox",
                "runs": runs,
            })
        return elements

    # SmartArt / other complex shapes: fall back to XML
    try:
        xml_elem = shape._element
        a_t_nodes = xml_elem.findall(".//" + "{" + NSMAP["a"] + "}t")
        if a_t_nodes:
            runs = []
            for i, node in enumerate(a_t_nodes):
                text = node.text
                if text and text.strip():
                    runs.append({"index": i, "text": text})
            if runs:
                counter["idx"] += 1
                elements.append({
                    "id": f"s{slide_num}_smartart_{counter['idx']}",
                    "type": "smartart",
                    "runs": runs,
                })
    except Exception:
        pass

    return elements


def extract_slides(pptx_path: str) -> list[dict]:
    """Extract all text elements from a PPTX file.

    Returns a list of slide dicts, each containing slide_number and elements.
    """
    prs = Presentation(pptx_path)
    slides_data = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        counter = {"idx": 0}
        elements = []

        # Process all shapes on the slide
        for shape in slide.shapes:
            elements.extend(_extract_shape(shape, slide_num, counter))

        # Process notes
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                runs = _extract_paragraphs(slide.notes_slide.notes_text_frame)
                if runs:
                    elements.append({
                        "id": f"s{slide_num}_note",
                        "type": "note",
                        "runs": runs,
                    })
        except Exception:
            pass

        slides_data.append({
            "slide_number": slide_num,
            "elements": elements,
        })

    return slides_data


# ── Apply helpers ──────────────────────────────────────────────────

def _build_translation_map(translations: list[dict]) -> dict:
    """Build a lookup: element_id -> {run_index: translated_text}."""
    tmap = {}
    for slide_data in translations:
        for elem in slide_data.get("elements", []):
            elem_id = elem.get("id")
            if not elem_id:
                continue
            run_map = {}
            for run in elem.get("runs", []):
                run_map[run["index"]] = run["text"]
            tmap[elem_id] = run_map
    return tmap


def _apply_to_text_frame(text_frame, elem_id: str, run_map: dict, stats: dict):
    """Apply translated runs to a text frame, preserving formatting."""
    global_idx = 0
    for para in text_frame.paragraphs:
        for run in para.runs:
            if not run.text:
                continue
            if global_idx in run_map:
                original_text = run.text
                translated_text = run_map[global_idx]
                try:
                    run.text = translated_text
                    # Change East Asian font to Arial
                    rPr = run._r.find("{" + NSMAP["a"] + "}rPr")
                    if rPr is not None:
                        ea = rPr.find("{" + NSMAP["a"] + "}ea")
                        if ea is not None:
                            ea.set("typeface", "Arial")
                        else:
                            etree.SubElement(rPr, "{" + NSMAP["a"] + "}ea").set("typeface", "Arial")
                    # AutoFit: scale font size proportionally
                    if run.font.size is not None and len(translated_text) > len(original_text):
                        ratio = len(original_text) / max(len(translated_text), 1)
                        new_size = max(run.font.size.pt * ratio, 8.0)
                        run.font.size = Pt(new_size)
                    stats["runs_ok"] += 1
                except Exception as e:
                    # Run-level fallback: keep original text
                    run.text = original_text
                    stats["runs_fail"] += 1
            global_idx += 1


def _apply_to_shape(shape, slide_num: int, counter: dict, tmap: dict, stats: dict):
    """Apply translations to a single shape, recursively handling groups."""
    # Group shapes
    if shape.shape_type is not None and shape.shape_type == 6:
        try:
            for child_shape in shape.shapes:
                _apply_to_shape(child_shape, slide_num, counter, tmap, stats)
        except Exception:
            pass
        return

    # Table
    if shape.has_table:
        table = shape.table
        for row_idx, row in enumerate(table.rows):
            for col_idx, cell in enumerate(row.cells):
                if cell.text_frame:
                    has_text = any(r.text for p in cell.text_frame.paragraphs for r in p.runs)
                    if has_text:
                        counter["idx"] += 1
                        elem_id = f"s{slide_num}_table_{counter['idx']}_r{row_idx}c{col_idx}"
                        if elem_id in tmap:
                            _apply_to_text_frame(cell.text_frame, elem_id, tmap[elem_id], stats)
                            stats["elements_ok"] += 1
                        else:
                            stats["elements_skip"] += 1
        return

    # Text frame
    if shape.has_text_frame:
        has_text = any(r.text for p in shape.text_frame.paragraphs for r in p.runs)
        if has_text:
            counter["idx"] += 1
            elem_id = f"s{slide_num}_shape_{counter['idx']}"
            if elem_id in tmap:
                _apply_to_text_frame(shape.text_frame, elem_id, tmap[elem_id], stats)
                stats["elements_ok"] += 1
            else:
                stats["elements_skip"] += 1
        return

    # SmartArt / XML fallback
    try:
        xml_elem = shape._element
        a_t_nodes = xml_elem.findall(".//" + "{" + NSMAP["a"] + "}t")
        text_nodes = [(i, n) for i, n in enumerate(a_t_nodes) if n.text and n.text.strip()]
        if text_nodes:
            counter["idx"] += 1
            elem_id = f"s{slide_num}_smartart_{counter['idx']}"
            if elem_id in tmap:
                run_map = tmap[elem_id]
                for i, node in text_nodes:
                    if i in run_map:
                        node.text = run_map[i]
                        stats["runs_ok"] += 1
                stats["elements_ok"] += 1
            else:
                stats["elements_skip"] += 1
    except Exception:
        pass


def apply_translations(
    pptx_path: str,
    translations: list[dict],
    output_path: str,
) -> dict:
    """Apply translated text back to PPTX, preserving formatting.

    Returns a stats dict with success/failure counts.
    """
    prs = Presentation(pptx_path)
    tmap = _build_translation_map(translations)
    stats = {
        "runs_ok": 0,
        "runs_fail": 0,
        "elements_ok": 0,
        "elements_skip": 0,
        "slides_ok": 0,
        "slides_fail": 0,
    }

    for slide_num, slide in enumerate(prs.slides, start=1):
        try:
            counter = {"idx": 0}
            for shape in slide.shapes:
                _apply_to_shape(shape, slide_num, counter, tmap, stats)

            # Notes
            try:
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    note_id = f"s{slide_num}_note"
                    if note_id in tmap:
                        _apply_to_text_frame(
                            slide.notes_slide.notes_text_frame,
                            note_id,
                            tmap[note_id],
                            stats,
                        )
                        stats["elements_ok"] += 1
            except Exception:
                pass

            stats["slides_ok"] += 1
        except Exception:
            stats["slides_fail"] += 1

    prs.save(output_path)
    return stats
